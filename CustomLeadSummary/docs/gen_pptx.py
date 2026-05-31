"""Generate the Web Lead Qualification PCF presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
BG_DARK    = RGBColor(0x1B, 0x1B, 0x1B)
BG_CARD    = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT     = RGBColor(0x00, 0x78, 0xD4)
ACCENT2    = RGBColor(0x10, 0x7C, 0x10)
ACCENT3    = RGBColor(0xFF, 0x8C, 0x00)
ACCENT4    = RGBColor(0xE7, 0x4C, 0x3C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY   = RGBColor(0x88, 0x88, 0x88)
TEAL       = RGBColor(0x00, 0xB2, 0x94)
PURPLE     = RGBColor(0x88, 0x17, 0xAF)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height,
              fill_color=None, line_color=None, line_width=Pt(0)):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = line_width
    else:
        shp.line.fill.background()
    return shp


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, color=WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items,
                    font_size=16, color=WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        run1 = p.add_run()
        run1.text = "\u25cf  "
        run1.font.size = Pt(font_size - 2)
        run1.font.color.rgb = bullet_color
        run1.font.name = "Segoe UI"
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = "Segoe UI"
    return txBox


def add_card(slide, left, top, width, height, title, body_items,
             accent_color=ACCENT, title_size=18, body_size=14):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                     left, top, width, height, fill_color=BG_CARD)
    card.adjustments[0] = 0.05
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, Pt(4),
              fill_color=accent_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.4),
                 title, font_size=title_size, color=WHITE, bold=True)
    if body_items:
        add_bullet_text(slide, left + Inches(0.2), top + Inches(0.55),
                        width - Inches(0.4), height - Inches(0.7),
                        body_items, font_size=body_size,
                        color=LIGHT_GRAY, bullet_color=accent_color)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0), Inches(0), Inches(0.15), Inches(7.5),
          fill_color=ACCENT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(1.2),
             "Web Lead Qualification",
             font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.7), Inches(10), Inches(0.8),
             "AI-Powered Lead Assessment for Dynamics 365 Sales",
             font_size=28, color=ACCENT)
add_text_box(slide, Inches(1), Inches(4.2), Inches(10), Inches(0.6),
             "PCF Control  \u00b7  Copilot Studio  \u00b7  MSAL Authentication  \u00b7  Power Automate",
             font_size=18, color=MED_GRAY)

badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  Inches(1), Inches(5.5), Inches(3.8), Inches(0.6),
                  fill_color=RGBColor(0x24, 0x29, 0x2E))
badge.adjustments[0] = 0.3
tf = badge.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
r.text = "Built entirely with GitHub Copilot"
r.font.size = Pt(16)
r.font.color.rgb = ACCENT2
r.font.bold = True
r.font.name = "Segoe UI"


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — The Challenge
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "The Challenge", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0.8), Inches(1.05), Inches(1.5), Pt(3), fill_color=ACCENT)

pain_data = [
    ("Slow Triage",
     "Manual review of every\nweb lead is slow\nand inconsistent", ACCENT4),
    ("Inconsistent",
     "Qualification varies by\nregion and individual\nsales rep", ACCENT3),
    ("Too Much Noise",
     "Spam, job seekers,\ncomplaints all land\nin the CRM", PURPLE),
    ("Missed Signals",
     "Text cues, domains,\nand intent are not\nevaluated consistently", ACCENT),
]

for i, (title, desc, color) in enumerate(pain_data):
    x = Inches(0.8 + i * 3.1)
    y = Inches(1.8)
    w = Inches(2.8)
    h = Inches(2.6)
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                     x, y, w, h, fill_color=BG_CARD)
    card.adjustments[0] = 0.06
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, Pt(4), fill_color=color)
    # Icon circle
    cx = x + (w - Inches(0.7)) / 2
    add_shape(slide, MSO_SHAPE.OVAL, cx, y + Inches(0.3),
              Inches(0.7), Inches(0.7), fill_color=color)
    add_text_box(slide, x, y + Inches(1.2), w, Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.7),
                 w - Inches(0.3), Inches(1.0),
                 desc, font_size=14, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

# Goal box
goal_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                     Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8),
                     fill_color=RGBColor(0x00, 0x3D, 0x6B))
goal_box.adjustments[0] = 0.1
add_text_box(slide, Inches(1.2), Inches(5.2), Inches(2), Inches(0.4),
             "GOAL", font_size=22, color=ACCENT, bold=True)
add_text_box(slide, Inches(1.2), Inches(5.7), Inches(10.8), Inches(0.9),
             "Provide an AI-assisted signal directly in the Lead form\n"
             "Advisory only  \u00b7  Explainable  \u00b7  Non-intrusive  \u00b7  Human-in-the-loop",
             font_size=18, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Design Principles
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Design Principles", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0.8), Inches(1.05), Inches(1.5), Pt(3), fill_color=ACCENT)

principles = [
    ("Advisory Only",
     "No automatic qualification\nchanges. AI suggests,\nhuman decides.", ACCENT),
    ("Explainable",
     "Summary + key factors\n+ disclaimer. Seller\nunderstands the why.", TEAL),
    ("Non-Intrusive",
     "Graceful errors, no form\nblocking. Assessment\nreadable in 5 seconds.", ACCENT2),
    ("Human-in-the-Loop",
     "Seller judgment is always\nauthoritative. Override\nanytime.", ACCENT3),
]

for i, (title, desc, color) in enumerate(principles):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.0)
    w = Inches(2.8)
    # Circle
    cx = x + (w - Inches(1.2)) / 2
    circle = add_shape(slide, MSO_SHAPE.OVAL, cx, y,
                       Inches(1.2), Inches(1.2), fill_color=color)
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(i + 1)
    r.font.size = Pt(32)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = "Segoe UI"
    # Title
    add_text_box(slide, x, y + Inches(1.5), w, Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Desc
    add_text_box(slide, x, y + Inches(2.1), w, Inches(1.2),
                 desc, font_size=15, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — Development Journey (Timeline)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Development Journey with GitHub Copilot",
             font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0.8), Inches(1.05), Inches(1.5), Pt(3), fill_color=ACCENT)

# Timeline bar
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(1.5), Inches(3.5), Inches(10.5), Pt(3),
          fill_color=MED_GRAY)

phases = [
    ("Research\n& Spec",     "Problem analysis\nUser stories\nAI schema design",  ACCENT),
    ("Code\nGeneration",    "7 TypeScript files\nReact components\nService layer",  ACCENT2),
    ("Testing &\nDev Harness", "22 unit tests\n6 mock scenarios\nLocal test server", TEAL),
    ("Deploy &\nPackage",   "Solution project\nManifest config\nProduction build",  ACCENT3),
    ("Production\nDebugging", "Architecture pivot\nMSAL integration\nEntity fixes",  ACCENT4),
]

for i, (title, desc, color) in enumerate(phases):
    x = Inches(1.0 + i * 2.3)
    # Dot on timeline
    dot = add_shape(slide, MSO_SHAPE.OVAL,
                    x + Inches(0.5), Inches(3.2),
                    Inches(0.6), Inches(0.6), fill_color=color)
    tf = dot.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(i + 1)
    r.font.size = Pt(18)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = "Segoe UI"
    # Title above
    add_text_box(slide, x, Inches(1.8), Inches(1.7), Inches(1.2),
                 title, font_size=16, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Desc below
    add_text_box(slide, x - Inches(0.2), Inches(4.1),
                 Inches(2.1), Inches(1.5),
                 desc, font_size=13, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

# Stat bar
stat_bar = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                     Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9),
                     fill_color=RGBColor(0x24, 0x29, 0x2E))
stat_bar.adjustments[0] = 0.2
stats = [("7", "Source Files"), ("22", "Unit Tests"),
         ("6", "Mock Scenarios"), ("7", "Iterations"), ("~2 MB", "Bundle")]
for i, (num, label) in enumerate(stats):
    sx = Inches(1.2 + i * 2.3)
    add_text_box(slide, sx, Inches(6.05), Inches(1.5), Inches(0.4),
                 num, font_size=24, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, sx, Inches(6.45), Inches(1.5), Inches(0.3),
                 label, font_size=12, color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — Solution Architecture
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Solution Architecture", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0.8), Inches(1.05), Inches(1.5), Pt(3), fill_color=ACCENT)

# D365 outer box
d365 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                 Inches(0.5), Inches(1.5), Inches(5.0), Inches(5.2),
                 line_color=ACCENT, line_width=Pt(2))
d365.adjustments[0] = 0.03
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(4), Inches(0.4),
             "Dynamics 365 Sales", font_size=18, color=ACCENT, bold=True)

# PCF box
pcf = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(2.2), Inches(4.4), Inches(3.8),
                fill_color=BG_CARD)
pcf.adjustments[0] = 0.04
add_text_box(slide, Inches(1.0), Inches(2.3), Inches(4), Inches(0.4),
             "PCF Control (React + Fluent UI)",
             font_size=16, color=WHITE, bold=True)

sub_cards = [
    ("Summary Card", "AI assessment + Key factors", ACCENT),
    ("Relevance Panel", "Toggle switch + Recommendations", TEAL),
    ("Feedback Panel", "Thumbs up/down + Comments", ACCENT3),
]
for i, (t, d, c) in enumerate(sub_cards):
    sy = Inches(2.9 + i * 1.05)
    sc = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   Inches(1.0), sy, Inches(4.0), Inches(0.9),
                   line_color=c, line_width=Pt(1.5))
    sc.adjustments[0] = 0.1
    add_text_box(slide, Inches(1.2), sy + Inches(0.1),
                 Inches(1.5), Inches(0.35),
                 t, font_size=13, color=c, bold=True)
    add_text_box(slide, Inches(2.8), sy + Inches(0.1),
                 Inches(2.0), Inches(0.7),
                 d, font_size=12, color=LIGHT_GRAY)

# External services
services = [
    (Inches(7.0), Inches(1.5), "Entra ID\n(MSAL)", PURPLE),
    (Inches(9.5), Inches(1.5), "Copilot\nStudio", ACCENT),
    (Inches(7.0), Inches(4.5), "Dataverse\nWeb API", TEAL),
    (Inches(9.5), Inches(3.2), "Power\nAutomate", ACCENT3),
    (Inches(11.5), Inches(3.2), "AI\nBuilder", ACCENT2),
]
for (x, y, label, color) in services:
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(1.8), Inches(1.3), fill_color=BG_CARD)
    box.adjustments[0] = 0.08
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              x, y, Inches(1.8), Pt(3), fill_color=color)
    add_text_box(slide, x, y + Inches(0.3), Inches(1.8), Inches(0.8),
                 label, font_size=14, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

# Arrows
add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
          Inches(5.3), Inches(2.0), Inches(1.5), Inches(0.25),
          fill_color=PURPLE)
add_text_box(slide, Inches(5.5), Inches(1.65), Inches(1.2), Inches(0.3),
             "Token", font_size=10, color=PURPLE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
          Inches(5.3), Inches(2.6), Inches(4.0), Inches(0.2),
          fill_color=ACCENT)
add_text_box(slide, Inches(6.5), Inches(2.3), Inches(2.0), Inches(0.3),
             "Event Activity (SSE)", font_size=10, color=ACCENT,
             bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
          Inches(5.3), Inches(5.0), Inches(1.5), Inches(0.2),
          fill_color=TEAL)
add_text_box(slide, Inches(5.5), Inches(4.7), Inches(1.3), Inches(0.3),
             "Feedback", font_size=10, color=TEAL, bold=True,
             alignment=PP_ALIGN.CENTER)

add_shape(slide, MSO_SHAPE.DOWN_ARROW,
          Inches(10.1), Inches(3.0), Inches(0.2), Inches(0.4),
          fill_color=ACCENT)
add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
          Inches(11.1), Inches(3.7), Inches(0.5), Inches(0.2),
          fill_color=ACCENT3)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — Authentication Flow (MSAL)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Authentication \u2014 MSAL Flow",
             font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0.8), Inches(1.05), Inches(1.5), Pt(3), fill_color=ACCENT)

# Three column headers
cols = [
    ("PCF Control", ACCENT, Inches(1.5)),
    ("Entra ID", PURPLE, Inches(5.5)),
    ("Copilot Studio API", TEAL, Inches(9.5)),
]
for (label, color, x) in cols:
    hdr = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, Inches(1.5), Inches(2.5), Inches(0.7),
                    fill_color=color)
    hdr.adjustments[0] = 0.2
    tf = hdr.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = label
    r.font.size = Pt(16)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = "Segoe UI"
    # Vertical line
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              x + Inches(1.22), Inches(2.2),
              Pt(2), Inches(4.2),
              fill_color=RGBColor(0x44, 0x44, 0x44))

# Steps
steps_data = [
    (Inches(2.5), "1. Check cached token", ACCENT),
    (Inches(3.1), "2. Token expired? Pop-up login", PURPLE),
    (Inches(3.7), "3. User authenticates (SSO / MFA)", ACCENT2),
    (Inches(4.3), "4. Access token received", ACCENT2),
    (Inches(4.9), "5. Bearer token + Event Activity", TEAL),
    (Inches(5.5), "6. Streaming response (SSE)", TEAL),
]
for (y, text, color) in steps_data:
    step_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(1.0), y, Inches(5.5), Inches(0.5),
                         line_color=color, line_width=Pt(1))
    step_box.adjustments[0] = 0.15
    add_text_box(slide, Inches(1.2), y + Inches(0.05),
                 Inches(5.0), Inches(0.4),
                 text, font_size=14, color=LIGHT_GRAY)

# Config card
add_card(slide, Inches(8.5), Inches(2.5), Inches(4.0), Inches(4.0),
         "Entra App Registration",
         [
             "App type: SPA",
             "Redirect: org.crm.dynamics.com",
             "Permission: CopilotStudio.Copilots.Invoke",
             "Scope: api.powerplatform.com/.default",
             "Admin consent granted",
             "Auth: Silent then Popup fallback",
         ], accent_color=PURPLE, body_size=13)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — Copilot Studio Agent
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Copilot Studio Agent Integration",
             font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0.8), Inches(1.05), Inches(1.5), Pt(3), fill_color=ACCENT)

# Code box
code_bg = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3),
                    fill_color=RGBColor(0x1E, 0x1E, 0x1E))
code_bg.adjustments[0] = 0.03
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.35),
             "Communication Protocol", font_size=14, color=ACCENT, bold=True)

code_lines = [
    ("// 1. Initialize SDK client", MED_GRAY),
    ("const client = new CopilotStudioClient(", WHITE),
    ("    { environmentId, agentIdentifier },", LIGHT_GRAY),
    ("    accessToken", LIGHT_GRAY),
    (");", WHITE),
    ("", WHITE),
    ("// 2. Start streaming conversation", MED_GRAY),
    ("for await (const a of", WHITE),
    ("    client.startConversationStreaming()) {", WHITE),
    ("    conversationId = a.conversation?.id;", LIGHT_GRAY),
    ("}", WHITE),
    ("", WHITE),
    ("// 3. Send Event Activity", MED_GRAY),
    ('const event = new Activity("event");', WHITE),
    ('event.name = "LeadQualification";', ACCENT2),
    ("event.value = { leadId };", ACCENT2),
    ("", WHITE),
    ("// 4. Receive streaming response", MED_GRAY),
    ("for await (const reply of", WHITE),
    ("    client.sendActivityStreaming(event)) {", WHITE),
    ("    agentResponseText = reply.text;", LIGHT_GRAY),
    ("}", WHITE),
]

y_off = Inches(2.1)
for (line, clr) in code_lines:
    add_text_box(slide, Inches(0.9), y_off, Inches(5.3), Inches(0.22),
                 line, font_size=11, color=clr, font_name="Consolas")
    y_off += Inches(0.22)

# Right cards
add_card(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.5),
         "Agent Topic Configuration",
         [
             'Trigger: "When an activity is received"',
             "Activity type: Event",
             "Event name: LeadQualification",
             "Input: System.Activity.Value.leadId",
             "Output: JSON message",
         ], accent_color=ACCENT, body_size=14)

add_card(slide, Inches(7.0), Inches(4.3), Inches(5.5), Inches(2.7),
         "AI Response Schema",
         [
             "summary \u2192 Natural language (1-3 sentences)",
             'isSalesRelevant \u2192 "yes" | "no" | "unclear"',
             "keyFactors \u2192 String array of reasons",
             "recommendationText \u2192 Explanation text",
             "fieldRecommendations \u2192 Field updates",
         ], accent_color=ACCENT2, body_size=14)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — End-to-End Flow
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "End-to-End Flow", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0.8), Inches(1.05), Inches(1.5), Pt(3), fill_color=ACCENT)

# Flow boxes — top row (left to right)
top_row = [
    ("User opens\nLead record",       ACCENT),
    ("PCF loads\nupdateView()",        RGBColor(0x55, 0x55, 0x55)),
    ("MSAL\nAcquire Token",           PURPLE),
    ("Start SSE\nConversation",        TEAL),
    ("Send Event\nLeadQualification",  ACCENT),
    ("Agent Topic\nTriggered",         ACCENT3),
]
# Bottom row (right to left)
bot_row = [
    ("Power Automate\nFetch Lead Data",   ACCENT3),
    ("AI Builder\nGenerate Assessment",   ACCENT2),
    ("Agent Returns\nJSON Response",      TEAL),
    ("PCF Parses &\nRenders UI",          ACCENT),
    ("User Reviews\n& Decides",           ACCENT2),
    ("Feedback saved\nto Dataverse",      PURPLE),
]

for i, (label, color) in enumerate(top_row):
    x = Inches(0.3 + i * 2.1)
    y = Inches(1.8)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(1.8), Inches(1.3), fill_color=BG_CARD)
    box.adjustments[0] = 0.08
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              x, y, Inches(1.8), Pt(3), fill_color=color)
    add_text_box(slide, x, y + Inches(0.25), Inches(1.8), Inches(0.9),
                 label, font_size=13, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

for i, (label, color) in enumerate(bot_row):
    x = Inches(0.3 + (5 - i) * 2.1)
    y = Inches(4.5)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(1.8), Inches(1.3), fill_color=BG_CARD)
    box.adjustments[0] = 0.08
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              x, y, Inches(1.8), Pt(3), fill_color=color)
    add_text_box(slide, x, y + Inches(0.25), Inches(1.8), Inches(0.9),
                 label, font_size=13, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

# Top row arrows
for i in range(5):
    ax = Inches(2.1 + i * 2.1)
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
              ax, Inches(2.3), Inches(0.35), Inches(0.18),
              fill_color=MED_GRAY)

# Down arrow right side
add_shape(slide, MSO_SHAPE.DOWN_ARROW,
          Inches(11.0), Inches(3.2), Inches(0.18), Inches(1.1),
          fill_color=MED_GRAY)

# Bottom row arrows (left-pointing)
for i in range(5):
    ax = Inches(10.1 - i * 2.1)
    add_shape(slide, MSO_SHAPE.LEFT_ARROW,
              ax, Inches(5.0), Inches(0.35), Inches(0.18),
              fill_color=MED_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — Lessons Learned
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Lessons Learned", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0.8), Inches(1.05), Inches(1.5), Pt(3), fill_color=ACCENT)

# Pivot highlight
pivot = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
                  fill_color=RGBColor(0x4A, 0x15, 0x15))
pivot.adjustments[0] = 0.06
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(11), Inches(0.4),
             "Architecture Pivot", font_size=20, color=ACCENT4, bold=True)
add_text_box(slide, Inches(0.9), Inches(2.1), Inches(5), Inches(0.5),
             "X  context.copilot.executeEvent()\n"
             '   "Result is not PredictAPI" in production',
             font_size=14, color=LIGHT_GRAY)
add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
          Inches(6.0), Inches(2.2), Inches(0.8), Inches(0.25),
          fill_color=ACCENT3)
add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(0.5),
             "MSAL + agents-copilotstudio-client SDK\n"
             "Microsoft-recommended reference architecture",
             font_size=14, color=ACCENT2)

# Gotcha cards
gotchas = [
    ("PCF Manifest", [
        "WebAPI feature must be declared",
        "Need bound property for form editor",
        "external-service-usage for HTTP",
    ], ACCENT),
    ("OData / Dataverse", [
        "Navigation names != column names",
        "Query relationship metadata",
        "Boolean vs OptionSet types",
    ], TEAL),
    ("Deployment", [
        "Production build (no eval)",
        "Version bump for cache bust",
        "Clear site data in browser",
    ], ACCENT3),
    ("Authentication", [
        "Scope != permission name",
        "Use .default suffix on scope",
        "SPA redirect URI required",
    ], PURPLE),
]
for i, (title, items, color) in enumerate(gotchas):
    x = Inches(0.5 + i * 3.15)
    add_card(slide, x, Inches(3.4), Inches(2.9), Inches(3.3),
             title, items, accent_color=color, body_size=13)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — Tech Stack + Thank You
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(0), Inches(0), Inches(0.15), Inches(7.5),
          fill_color=ACCENT)

add_text_box(slide, Inches(1), Inches(0.8), Inches(10), Inches(0.8),
             "Technology Stack", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE,
          Inches(1), Inches(1.45), Inches(1.5), Pt(3), fill_color=ACCENT)

tech_items = [
    ("React 16",           "UI Framework",    ACCENT),
    ("Fluent UI v9",       "Design System",   TEAL),
    ("TypeScript",         "Language",         ACCENT2),
    ("MSAL v4",            "Authentication",   PURPLE),
    ("CopilotStudio SDK",  "Agent Client",     ACCENT),
    ("Jest",               "Testing",          ACCENT3),
    ("Webpack",            "Build Tool",       ACCENT4),
    ("Dataverse",          "Data Platform",    TEAL),
    ("AI Builder",         "AI Engine",        ACCENT2),
    ("GitHub Copilot",     "Dev Partner",      RGBColor(0xF0, 0xC0, 0x00)),
]

for i, (name, desc, color) in enumerate(tech_items):
    row = i // 5
    col = i % 5
    x = Inches(0.8 + col * 2.5)
    y = Inches(2.0 + row * 1.5)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                    x, y, Inches(2.2), Inches(1.1), fill_color=BG_CARD)
    box.adjustments[0] = 0.08
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              x, y, Pt(4), Inches(1.1), fill_color=color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15),
                 Inches(2.0), Inches(0.4),
                 name, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55),
                 Inches(2.0), Inches(0.3),
                 desc, font_size=12, color=MED_GRAY)

# Thank you
add_text_box(slide, Inches(1), Inches(5.5), Inches(10), Inches(0.6),
             "Thank You", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             "Reference Architecture: learn.microsoft.com/dynamics365/guidance/"
             "reference-architectures/custom-copilot-agent-dynamics-365-power-apps",
             font_size=13, color=ACCENT)


# ═══════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════
out = os.path.join(
    r"C:\Users\petkra\Source\CustomLeadSummary\docs",
    "Web Lead Qualification PCF.pptx",
)
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
